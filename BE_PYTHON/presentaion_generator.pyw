from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import image_generator
import re
import ast
import json
from datetime import date


def genrater(data,Head):
    print(Head)
        # Your data
    data1 = '''
    Titles : [ESOFT Metro Campus, Private Sector Education, Computer Studies, Academic and Professional], Notes : [ESOFT Metro Campus is a private sector educational institute/college which was formerly known as ESOFT Computer Studies. It is based in Colombo, Sri Lanka and has been offering academic and professional qualifications in Computing, Business, Management, Engineering, Hospitality and English since its establishment in the year. This company is headed by Dr Dayan Rajapakse and presently has more than branches all across the nation.]'''
   # Extract the titles and notes using string manipulation
    titles_start = data.find('Titles : [')
    notes_start = data.find('Notes : [')

    titles = data[titles_start:notes_start].strip()
    notes = data[notes_start:].strip()

    # Split the titles and notes into lists
    titles = titles.split(' : [')[1].strip().strip(']').split(', ')
    notes = notes.split(' : [')[1].strip().strip(']').split('. ')
    print(titles)

    print('\n',notes)
    # Create a list of dictionaries
    data_list = [{'Title': title, 'Notes': note} for title, note in zip(titles, notes)]

    print (data_list)

    # Open the PowerPoint presentation
    prs = Presentation('Template.pptx')

  
    # Download images for the specified Title
    print("Head is"+Head)
    #image_generator.download_images(Head)

    # Update the title of the first slide

    # Open the file in read mode
    with open('User/mail.txt', 'r') as file:
        # Read the content of the file
        creater_name = file.read()
        
    #Title='Bridging the Accessibility Gap with AI Voice Assistants'
    prs.slides[0].placeholders[0].text = Head
    prs.slides[0].placeholders[13].text = str(creater_name)
    # Set the path to the image
    image_path = 'STORE/' + Head + '/' + Head + '1.png'

    # Set the picture in Slide 1, Placeholder 21
    slide = prs.slides[0]
    picture = slide.placeholders[21]
 
    # Replace the existing picture with the new one
    picture.insert_picture(image_path)

    image_path = 'STORE/' + Head + '/' + Head + '2.png'

    # Remove unwanted symbols from titles and join them into a content list
    unwanted_symbols = ''']."','''
    content_list = '\n'.join(entry['Title'].translate(str.maketrans('', '', unwanted_symbols)) for entry in data_list)

    # Update the content list on the second slide
    prs.slides[1].placeholders[14].text = content_list
    # Set the picture in Slide 1, Placeholder 21
    slide = prs.slides[1]
    picture = slide.placeholders[15]
    # Replace the existing picture with the new one
    picture.insert_picture(image_path)# Loop through the data and update the slides

    slide = prs.slides[9]
    picture = slide.placeholders[21]
        # Replace the existing picture with the new one
    image_path3 = 'STORE/' + Head + '/' + Head+'6.png'
    picture.insert_picture(image_path3)# Loop through the data and update the slides

    
    image_path2 = 'STORE/' + Head + '/' + Head 
    unwanted_symbols = '].",'
    num=2
    for slide_number, entry in enumerate(data_list, start=2):
        slide = prs.slides[slide_number]
        image_path=image_path2+str(num)+'.png'
        # Update the title and notes
        updated_title = entry['Title'].translate(str.maketrans('', '', unwanted_symbols))
        updated_notes = entry['Notes'].translate(str.maketrans('', '', unwanted_symbols))

        slide.placeholders[0].text = updated_title
        slide.placeholders[17].text = updated_notes
        picture = slide.placeholders[18]
        picture.insert_picture(image_path)
        num=num+1
        # Format the text
        text_frame = slide.placeholders[17].text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(26)
                run.font.color.rgb = RGBColor(0, 0, 0)
                run.font.bold = False  # Set the text to not bold
    # Save the updated presentation in the desired location
         # Get today's date
        today = date.today()
        # Format the date as a string
        formatted_date = today.strftime("%Y-%m-%d")

        prs.slides[9].placeholders[20].text = formatted_date

        
    presentation_path = 'STORE/' + Head + '/' + Head + '.pptx'
    prs.save(presentation_path)

#genrater( 'Bridging the Accessibility Gap with AI Voice Assistants')
