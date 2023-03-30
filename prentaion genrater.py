from pptx import Presentation
import requests
from io import BytesIO

# open the PowerPoint presentation
prs = Presentation('PRESENTATION1.pptx')

# get the slide you want to modify (slide 256 in this case)
slide = prs.slides[0]

# access the placeholders and update the text
title = slide.placeholders[0]
title.text = 'New Title'

text1 = slide.placeholders[13]
text1.text = 'New Text for Placeholder 3'


# set an online image for Placeholder 21
image_url = 'https://images.idgesg.net/images/article/2017/10/wireless_network_internet_of_things_iot_thinkstock_853701554_3x2-100740688-large.jpg?auto=webp&quality=85,70'
image_request = requests.get(image_url)
image_content = BytesIO(image_request.content)
picture = slide.placeholders[21]
picture.insert_picture(image_content)


#text2 = slide.placeholders[20]
#text2.text = 'New Text for Placeholder 2'

 #save the updated presentation
prs.save('PRESENTATION1_updated.pptx')

# get placeholders and their names on the slide
for placeholder in slide.placeholders:
    print(f"Placeholder {placeholder.placeholder_format.idx}: {placeholder.name}")
